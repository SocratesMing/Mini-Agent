import os
import tempfile
from typing import Optional, Dict, Any, List

from mini_agent.tools.base import Tool, ToolResult

CODE_FILE_EXTENSIONS = {
    'py': 'Python',
    'js': 'JavaScript',
    'ts': 'TypeScript',
    'jsx': 'React JSX',
    'tsx': 'React TSX',
    'java': 'Java',
    'c': 'C',
    'cpp': 'C++',
    'cc': 'C++',
    'cxx': 'C++',
    'h': 'C/C++ Header',
    'hpp': 'C++ Header',
    'cs': 'C#',
    'go': 'Go',
    'rs': 'Rust',
    'rb': 'Ruby',
    'php': 'PHP',
    'swift': 'Swift',
    'kt': 'Kotlin',
    'scala': 'Scala',
    'r': 'R',
    'm': 'Objective-C',
    'mm': 'Objective-C++',
    'sh': 'Shell',
    'bash': 'Bash',
    'zsh': 'Zsh',
    'ps1': 'PowerShell',
    'bat': 'Batch',
    'sql': 'SQL',
    'html': 'HTML',
    'htm': 'HTML',
    'css': 'CSS',
    'scss': 'SCSS',
    'sass': 'Sass',
    'less': 'Less',
    'xml': 'XML',
    'json': 'JSON',
    'yaml': 'YAML',
    'yml': 'YAML',
    'toml': 'TOML',
    'ini': 'INI',
    'cfg': 'Config',
    'conf': 'Config',
    'md': 'Markdown',
    'markdown': 'Markdown',
    'vue': 'Vue',
    'svelte': 'Svelte',
    'lua': 'Lua',
    'pl': 'Perl',
    'pm': 'Perl Module',
    'raku': 'Raku',
    'dart': 'Dart',
    'clj': 'Clojure',
    'ex': 'Elixir',
    'exs': 'Elixir Script',
    'erl': 'Erlang',
    'hrl': 'Erlang Header',
    'hs': 'Haskell',
    'lhs': 'Literate Haskell',
    'ml': 'OCaml',
    'mli': 'OCaml Interface',
    'fs': 'F#',
    'fsi': 'F# Interface',
    'v': 'Verilog',
    'sv': 'SystemVerilog',
    'vhdl': 'VHDL',
    'asm': 'Assembly',
    's': 'Assembly',
}

class DocumentParseTool(Tool):
    """解析各种格式的文档文件，提取文本内容"""
    
    @property
    def name(self) -> str:
        return "DocumentParseTool"
    
    @property
    def description(self) -> str:
        return "解析各种格式的文档文件，提取文本内容。支持 PDF、Word、Excel、PowerPoint、TXT、CSV、JSON、XML、Markdown 以及各种编程语言代码文件。"
    
    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "文件路径"
                }
            },
            "required": ["file_path"]
        }
    
    def __init__(self):
        self._ensure_dependencies()
    
    def _ensure_dependencies(self):
        """确保必要的依赖已安装"""
        try:
            import PyPDF2
        except ImportError:
            pass
        try:
            import docx
        except ImportError:
            pass
        try:
            import pandas
        except ImportError:
            pass
        try:
            import pptx
        except ImportError:
            pass
        try:
            import chardet
        except ImportError:
            pass
    
    def _detect_encoding(self, file_path: str) -> str:
        """检测文件编码"""
        try:
            import chardet
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8')
        except ImportError:
            return 'utf-8'
        except Exception:
            return 'utf-8'

    def _parse_txt(self, file_path: str) -> str:
        """解析文本文件，自动检测编码"""
        try:
            encoding = self._detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                return f.read()
        except Exception as e:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e2:
                return f"Error parsing TXT file: {str(e2)}"

    def _parse_pdf(self, file_path: str) -> str:
        """解析PDF文件"""
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)
                text += f"PDF 文档共 {total_pages} 页\n\n"
                for page_num in range(total_pages):
                    page = reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += f"--- 第 {page_num + 1} 页 ---\n"
                        text += page_text + "\n\n"
            return text.strip()
        except ImportError:
            return "PyPDF2 not installed, cannot parse PDF files"
        except Exception as e:
            return f"Error parsing PDF file: {str(e)}"

    def _parse_docx(self, file_path: str) -> str:
        """解析Word文档"""
        try:
            import docx
            doc = docx.Document(file_path)
            text_parts = []
            
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    for paragraph in doc.paragraphs:
                        if paragraph._element == element:
                            text_parts.append(paragraph.text)
                            break
                elif element.tag.endswith('tbl'):
                    for table in doc.tables:
                        if table._element == element:
                            table_text = self._parse_docx_table(table)
                            text_parts.append(table_text)
                            break
            
            return "\n".join(text_parts).strip()
        except ImportError:
            return "python-docx not installed, cannot parse DOCX files"
        except Exception as e:
            return f"Error parsing DOCX file: {str(e)}"
    
    def _parse_docx_table(self, table) -> str:
        """解析 Word 文档中的表格"""
        try:
            import docx
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            return "\n".join(rows)
        except Exception:
            return ""

    def _parse_excel(self, file_path: str) -> str:
        """解析Excel文件"""
        try:
            import pandas as pd
            text_parts = []
            xl = pd.ExcelFile(file_path)
            
            text_parts.append(f"Excel 文件共 {len(xl.sheet_names)} 个工作表: {', '.join(xl.sheet_names)}\n")
            
            for sheet_name in xl.sheet_names:
                text_parts.append(f"\n=== 工作表: {sheet_name} ===\n")
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                max_rows = 1000
                if len(df) > max_rows:
                    text_parts.append(f"(数据较多，仅显示前 {max_rows} 行)\n")
                    df = df.head(max_rows)
                
                text_parts.append(df.to_string(index=False))
                text_parts.append(f"\n共 {len(df)} 行数据")
            
            return "\n".join(text_parts)
        except ImportError:
            return "pandas not installed, cannot parse Excel files"
        except Exception as e:
            return f"Error parsing Excel file: {str(e)}"

    def _parse_csv(self, file_path: str) -> str:
        """解析CSV文件"""
        try:
            import pandas as pd
            encoding = self._detect_encoding(file_path)
            df = pd.read_csv(file_path, encoding=encoding)
            
            max_rows = 1000
            if len(df) > max_rows:
                result = f"(数据较多，仅显示前 {max_rows} 行)\n\n"
                result += df.head(max_rows).to_string(index=False)
                result += f"\n\n共 {len(df)} 行数据"
                return result
            
            return df.to_string(index=False)
        except ImportError:
            return "pandas not installed, cannot parse CSV files"
        except Exception as e:
            return f"Error parsing CSV file: {str(e)}"
    
    def _parse_pptx(self, file_path: str) -> str:
        """解析 PowerPoint 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            
            text_parts.append(f"PowerPoint 文档共 {len(prs.slides)} 张幻灯片\n")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n=== 幻灯片 {slide_num} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                    
                    if shape.has_table:
                        table_text = self._parse_pptx_table(shape.table)
                        if table_text:
                            text_parts.append(table_text)
            
            return "\n".join(text_parts)
        except ImportError:
            return "python-pptx not installed, cannot parse PPTX files"
        except Exception as e:
            return f"Error parsing PPTX file: {str(e)}"
    
    def _parse_pptx_table(self, table) -> str:
        """解析 PowerPoint 中的表格"""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            return "\n".join(rows)
        except Exception:
            return ""
    
    def _parse_json(self, file_path: str) -> str:
        """解析 JSON 文件"""
        try:
            import json
            encoding = self._detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as f:
                data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            return f"Error parsing JSON file: {str(e)}"
    
    def _parse_xml(self, file_path: str) -> str:
        """解析 XML 文件"""
        try:
            import xml.etree.ElementTree as ET
            encoding = self._detect_encoding(file_path)
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            def element_to_text(element, level=0):
                indent = "  " * level
                text_parts = []
                
                if element.text and element.text.strip():
                    text_parts.append(f"{indent}{element.tag}: {element.text.strip()}")
                
                for child in element:
                    text_parts.append(element_to_text(child, level + 1))
                
                return "\n".join(text_parts)
            
            return element_to_text(root)
        except Exception as e:
            return f"Error parsing XML file: {str(e)}"
    
    def _parse_code(self, file_path: str, language: str) -> str:
        """解析代码文件"""
        try:
            encoding = self._detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                code = f.read()
            
            lines = code.split('\n')
            result = f"代码文件 ({language})\n"
            result += f"文件路径: {file_path}\n"
            result += f"总行数: {len(lines)}\n\n"
            result += "--- 代码内容 ---\n\n"
            result += code
            
            return result
        except Exception as e:
            return f"Error parsing code file: {str(e)}"
    
    def _parse_markdown(self, file_path: str) -> str:
        """解析 Markdown 文件"""
        try:
            encoding = self._detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                content = f.read()
            return content
        except Exception as e:
            return f"Error parsing Markdown file: {str(e)}"

    def _detect_file_type(self, file_path: str) -> str:
        """检测文件类型"""
        _, ext = os.path.splitext(file_path.lower())
        return ext[1:] if ext else ''

    def run(self, file_path: str) -> Dict[str, Any]:
        """运行文档解析工具"""
        if not os.path.exists(file_path):
            return {
                "success": False,
                "content": f"File not found: {file_path}"
            }
        
        file_type = self._detect_file_type(file_path)
        file_info = {
            "file_path": file_path,
            "file_type": file_type,
            "file_name": os.path.basename(file_path),
        }
        
        try:
            content = ""
            
            if file_type == 'txt':
                content = self._parse_txt(file_path)
            elif file_type == 'pdf':
                content = self._parse_pdf(file_path)
            elif file_type == 'docx':
                content = self._parse_docx(file_path)
            elif file_type == 'doc':
                content = self._parse_docx(file_path)
            elif file_type in ['xlsx', 'xls']:
                content = self._parse_excel(file_path)
            elif file_type == 'csv':
                content = self._parse_csv(file_path)
            elif file_type in ['pptx', 'ppt']:
                content = self._parse_pptx(file_path)
            elif file_type == 'json':
                content = self._parse_json(file_path)
            elif file_type == 'xml':
                content = self._parse_xml(file_path)
            elif file_type in ['md', 'markdown']:
                content = self._parse_markdown(file_path)
            elif file_type in CODE_FILE_EXTENSIONS:
                language = CODE_FILE_EXTENSIONS[file_type]
                content = self._parse_code(file_path, language)
            else:
                try:
                    content = self._parse_txt(file_path)
                    if not content.strip() or len(content) < 10:
                        content = f"Unsupported file type: {file_type}"
                except Exception:
                    content = f"Unsupported file type: {file_type}"
            
            return {
                "success": True,
                **file_info,
                "content": content,
                "content_length": len(content)
            }
        except Exception as e:
            return {
                "success": False,
                **file_info,
                "content": f"Error parsing document: {str(e)}"
            }
    
    async def execute(self, file_path: str, **kwargs) -> ToolResult:
        """Execute the tool."""
        result = self.run(file_path)
        if result.get("success"):
            return ToolResult(
                success=True,
                content=result.get("content", "")
            )
        else:
            return ToolResult(
                success=False,
                content=result.get("content", ""),
                error=result.get("content", "Unknown error")
            )

class DocumentInfoTool(Tool):
    """获取文档信息"""
    
    @property
    def name(self) -> str:
        return "DocumentInfoTool"
    
    @property
    def description(self) -> str:
        return "获取文档基本信息"
    
    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "文件路径"
                }
            },
            "required": ["file_path"]
        }
    
    def run(self, file_path: str) -> Dict[str, Any]:
        """运行文档信息工具"""
        if not os.path.exists(file_path):
            return {
                "success": False,
                "info": f"File not found: {file_path}"
            }
        
        try:
            file_size = os.path.getsize(file_path)
            file_type = os.path.splitext(file_path.lower())[1][1:]
            
            # 格式化文件大小
            def format_size(size):
                if size < 1024:
                    return f"{size} B"
                elif size < 1024 * 1024:
                    return f"{size / 1024:.2f} KB"
                else:
                    return f"{size / (1024 * 1024):.2f} MB"
            
            info = {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_type": file_type,
                "file_size": format_size(file_size),
                "size_bytes": file_size,
                "exists": True
            }
            
            return {
                "success": True,
                "info": info
            }
        except Exception as e:
            return {
                "success": False,
                "info": f"Error getting document info: {str(e)}"
            }
    
    async def execute(self, file_path: str, **kwargs) -> ToolResult:
        """Execute the tool."""
        result = self.run(file_path)
        if result.get("success"):
            return ToolResult(
                success=True,
                content=str(result.get("info", {}))
            )
        else:
            return ToolResult(
                success=False,
                content=result.get("info", ""),
                error=result.get("info", "Unknown error")
            )